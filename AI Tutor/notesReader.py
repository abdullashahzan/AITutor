from pdfminer.high_level import extract_text
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2

def clean_text(text):
    # Remove newline characters
    text = text.replace('\n', ' ')
    # Remove extra spaces
    text = ' '.join(text.split())
    return text

def extract_text_from_pdf(pdf_file):
    text = extract_text(pdf_file)
    return clean_text(text)

def extract_text_from_ppt(ppt_file):
    # Load the PowerPoint file
    prs = Presentation(ppt_file)

    # Initialize an empty string to store all text
    text = ""

    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return clean_text(text)

def extract_text_from_txt(txt_file):
    with open(txt_file, 'r', encoding='utf-8') as file:
        text = file.read()
    return clean_text(text)

def extract_text_from_image(image_file):
    # Open the image file
    #img = Image.open(image_file)

    img = cv2.imread(image_file)  # Load the image
    img = cv2.resize(img, (0, 0), fx=0.25, fy=0.25)
    img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)  # convert to grey
    img = cv2.adaptiveThreshold(img, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, 3, 15)

    # Use pytesseract to extract text
    text = pytesseract.image_to_string(img)
    # Clean the extracted text
    return clean_text(text)

import os

def extract_text_from_file(file_path):
    # Get the file extension
    _, file_extension = os.path.splitext(file_path)
    
    # Determine the extraction type based on the file extension
    if file_extension in ['.pdf']:
        return extract_text_from_pdf(file_path)
    elif file_extension in ['.ppt', '.pptx']:
        return extract_text_from_ppt(file_path)
    elif file_extension in ['.txt']:
        return extract_text_from_txt(file_path)
    elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
        return extract_text_from_image(file_path)
    else:
        return 'Unsupported file type'
